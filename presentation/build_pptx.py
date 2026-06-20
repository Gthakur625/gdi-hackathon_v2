"""
JAGGU AI Hackathon Presentation Builder
Velocity Branding · Purple Gradient · Premium SaaS Design
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Brand Palette ─────────────────────────────────────────────────────────────
PURPLE_DARK   = RGBColor(0x0B, 0x0F, 0x19)   # #0B0F19
PURPLE_MID    = RGBColor(0x4F, 0x46, 0xE5)   # #4F46E5
PURPLE_BRIGHT = RGBColor(0x7C, 0x3A, 0xED)   # #7C3AED
PURPLE_LIGHT  = RGBColor(0xC0, 0x84, 0xFC)   # #C084FC
ACCENT_GREEN  = RGBColor(0x34, 0xD3, 0x99)   # #34D399
ACCENT_AMBER  = RGBColor(0xFB, 0xBF, 0x24)   # #FBBF24
ACCENT_RED    = RGBColor(0xF8, 0x71, 0x71)   # #F87171
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LIGHT    = RGBColor(0xD1, 0xD5, 0xDB)   # #D1D5DB
GREY_MID      = RGBColor(0x9C, 0xA3, 0xAF)   # #9CA3AF
INDIGO_LINE   = RGBColor(0x1F, 0x29, 0x37)   # #1F2937

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # fully blank


# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, fill_alpha=None, line_rgb=None, line_width=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width:
        shape.line.color.rgb = line_rgb
        shape.line.width     = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    if hasattr(run.font, 'name'):
        run.font.name = "Calibri"
    return txBox


def add_multiline(slide, lines, l, t, w, h,
                  size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                  line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        if hasattr(run.font, 'name'):
            run.font.name = "Calibri"
    return txBox


def gradient_bg(slide):
    """Dark navy background."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=PURPLE_DARK)


def accent_bar(slide, y=0.0, h=0.05, color=PURPLE_MID):
    """Thin horizontal accent bar."""
    add_rect(slide, 0, y, 13.33, h, fill_rgb=color)


def card(slide, l, t, w, h, fill=RGBColor(0x11, 0x18, 0x27), border=PURPLE_MID):
    """Dark card with purple border."""
    add_rect(slide, l, t, w, h, fill_rgb=fill, line_rgb=border, line_width=0.8)


def bullet_block(slide, items, l, t, w, icon="▸", size=13, color=GREY_LIGHT, icon_color=PURPLE_LIGHT):
    y = t
    for item in items:
        add_text(slide, icon, l, y, 0.3, 0.3, size=size, bold=True, color=icon_color)
        add_text(slide, item, l+0.28, y, w-0.28, 0.3, size=size, color=color)
        y += 0.32
    return y


def pill(slide, text, l, t, w=1.4, h=0.32, fill=PURPLE_MID, text_color=WHITE, size=11):
    add_rect(slide, l, t, w, h, fill_rgb=fill)
    add_text(slide, text, l, t+0.02, w, h, size=size, bold=True, color=text_color,
             align=PP_ALIGN.CENTER)


def velocity_watermark(slide):
    add_text(slide, "VELOCITY GDI", 0.3, 7.1, 3, 0.3,
             size=8, color=GREY_MID, bold=True)
    add_text(slide, "Confidential · Hackathon Demo 2026", 9.5, 7.1, 3.5, 0.3,
             size=8, color=GREY_MID, align=PP_ALIGN.RIGHT)


def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def arrow_flow(slide, steps, l, t, w_total, h=0.55, box_color=PURPLE_MID, arrow_color=PURPLE_LIGHT):
    n = len(steps)
    box_w = (w_total - (n-1)*0.25) / n
    x = l
    for i, step in enumerate(steps):
        card(slide, x, t, box_w, h, fill=box_color, border=PURPLE_LIGHT)
        add_text(slide, step, x, t, box_w, h, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < n-1:
            add_text(slide, "→", x+box_w, t+0.1, 0.25, 0.35, size=16, bold=True,
                     color=arrow_color, align=PP_ALIGN.CENTER)
            x += box_w + 0.25
        else:
            x += box_w


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
gradient_bg(s1)
accent_bar(s1, y=0.0, h=0.06, color=PURPLE_MID)
accent_bar(s1, y=7.44, h=0.06, color=PURPLE_BRIGHT)

# Slide label
add_text(s1, "01 / 05", 12.6, 0.12, 0.6, 0.3, size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)

# Title block
add_text(s1, "The Problem", 0.5, 0.2, 8, 0.8, size=38, bold=True, color=WHITE)
add_text(s1, "D2C Brands Lose Revenue Due to Failed Deliveries",
         0.5, 0.95, 9, 0.45, size=17, color=PURPLE_LIGHT, italic=True)

# Left column — challenges
card(s1, 0.4, 1.55, 3.8, 4.6, fill=RGBColor(0x11, 0x18, 0x27), border=PURPLE_MID)
add_text(s1, "Key Challenges", 0.55, 1.65, 3.5, 0.4, size=13, bold=True, color=ACCENT_AMBER)
bullet_block(s1, [
    "High RTO (Return to Origin)",
    "High NDR (Non-Delivery)",
    "Customer Unreachable",
    "Courier Underperformance",
    "Poor Courier Allocation",
    "Lack of Operational Visibility",
    "Manual Analysis & Reporting",
], l=0.55, t=2.05, w=3.5, size=12, icon="✗", icon_color=ACCENT_RED)

# Middle column — example impact
card(s1, 4.4, 1.55, 4.2, 4.6, fill=RGBColor(0x0F, 0x17, 0x26), border=ACCENT_RED)
add_text(s1, "Business Impact", 4.55, 1.65, 3.9, 0.4, size=13, bold=True, color=ACCENT_RED)
add_text(s1, "Example Seller", 4.55, 2.1, 3.9, 0.32, size=11, color=GREY_MID, bold=True)

# Big numbers
add_text(s1, "10,000", 4.55, 2.45, 3.9, 0.7, size=40, bold=True, color=WHITE)
add_text(s1, "Orders / Month", 4.55, 3.1, 3.9, 0.3, size=11, color=GREY_MID)

add_rect(s1, 4.55, 3.45, 3.9, 0.02, fill_rgb=INDIGO_LINE)

add_text(s1, "65%", 4.55, 3.55, 1.5, 0.6, size=32, bold=True, color=ACCENT_AMBER)
add_text(s1, "Delivery\nRate", 6.0, 3.6, 2, 0.6, size=11, color=GREY_MID)

add_text(s1, "3,500", 4.55, 4.25, 1.5, 0.6, size=32, bold=True, color=ACCENT_RED)
add_text(s1, "Orders\nFailed", 6.0, 4.3, 2, 0.6, size=11, color=GREY_MID)

add_text(s1, "Revenue at risk every month", 4.55, 4.95, 3.9, 0.35,
         size=11, color=GREY_MID, italic=True)

# Right — flow diagram
card(s1, 8.8, 1.55, 4.1, 4.6, fill=RGBColor(0x11, 0x18, 0x27), border=PURPLE_BRIGHT)
add_text(s1, "The Failure Chain", 8.95, 1.65, 3.8, 0.4, size=13, bold=True, color=PURPLE_LIGHT)

flow_items = [
    ("📦 Order Placed",    WHITE,        RGBColor(0x1F,0x29,0x37)),
    ("🚚 Dispatched",      GREY_LIGHT,   RGBColor(0x1F,0x29,0x37)),
    ("⚠️ NDR",             ACCENT_AMBER, RGBColor(0x2D,0x1F,0x07)),
    ("🔄 Reattempt",       GREY_LIGHT,   RGBColor(0x1F,0x29,0x37)),
    ("❌ RTO",             ACCENT_RED,   RGBColor(0x2D,0x07,0x07)),
    ("💸 Revenue Loss",    ACCENT_RED,   RGBColor(0x2D,0x07,0x07)),
]
fy = 2.15
for label, tc, fc in flow_items:
    add_rect(s1, 9.0, fy, 2.7, 0.36, fill_rgb=fc, line_rgb=INDIGO_LINE, line_width=0.5)
    add_text(s1, label, 9.05, fy+0.04, 2.6, 0.3, size=11, bold=True, color=tc)
    if label != "💸 Revenue Loss":
        add_text(s1, "↓", 9.0, fy+0.36, 2.7, 0.2, size=10, color=GREY_MID, align=PP_ALIGN.CENTER)
        fy += 0.57
    fy += 0.0

velocity_watermark(s1)
add_notes(s1, """SPEAKER NOTES — Slide 1: The Problem

Open strong: "Every D2C brand ships thousands of orders. But up to 35% never get delivered."

Key talking points:
- The average D2C brand on Velocity loses 25–35% of shipments to NDR and RTO
- For a seller doing 10,000 orders/month at ₹1,200 AOV, that's ₹4.2 crore in failed delivery revenue
- Root causes: no real-time intelligence, manual operations, courier blind spots
- Current tools are dashboards — they show what happened, not what to do

End: "This is the problem JAGGU AI was built to solve." """)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCING JAGGU AI
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
gradient_bg(s2)
accent_bar(s2, y=0.0, h=0.06, color=PURPLE_BRIGHT)
accent_bar(s2, y=7.44, h=0.06, color=PURPLE_MID)
add_text(s2, "02 / 05", 12.6, 0.12, 0.6, 0.3, size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)

# Gradient title area
add_rect(s2, 0, 0.06, 13.33, 1.5, fill_rgb=RGBColor(0x0F, 0x0A, 0x1E))
add_text(s2, "🤖  JAGGU AI", 0.4, 0.15, 12, 0.75, size=44, bold=True, color=WHITE)
add_text(s2, "Powered by Velocity GDI — Growth & Delivery Intelligence",
         0.4, 0.88, 10, 0.38, size=15, color=PURPLE_LIGHT, italic=True)

# Tagline pill
pill(s2, "Your AI KAM & Operations Expert", 0.4, 1.42, 4.8, 0.38,
     fill=PURPLE_MID, size=13)

# Not a dashboard / chatbot statement
card(s2, 0.4, 2.0, 5.2, 1.35, fill=RGBColor(0x1A, 0x0A, 0x2E), border=PURPLE_BRIGHT)
add_text(s2, "Not a Dashboard.  Not a Chatbot.", 0.6, 2.1, 4.8, 0.38,
         size=15, bold=True, color=ACCENT_AMBER)
add_text(s2, "An AI Growth & Delivery Consultant.", 0.6, 2.48, 4.8, 0.38,
         size=14, color=WHITE, italic=True)
add_text(s2, "JAGGU gives proactive, personalised insights — not reports.",
         0.6, 2.86, 4.8, 0.35, size=11, color=GREY_LIGHT)

# 6 capability cards
caps = [
    ("📊", "Delivery\nIntelligence",  PURPLE_MID),
    ("🚚", "Courier\nIntelligence",   PURPLE_BRIGHT),
    ("📦", "Product\nIntelligence",   RGBColor(0x06,0x95,0x6C)),
    ("📍", "Pincode\nIntelligence",   RGBColor(0x92,0x40,0x0E)),
    ("⚙️", "VAS\nRecommendations",   RGBColor(0x6D,0x28,0xD9)),
    ("💡", "Operational\nIntelligence",RGBColor(0x1D,0x4E,0xD8)),
]
cx = 0.4; cy = 3.5; cw = 2.05; ch = 1.1
for i, (icon, label, clr) in enumerate(caps):
    cc = cx + i*(cw+0.07)
    card(s2, cc, cy, cw, ch, fill=RGBColor(0x11,0x18,0x27), border=clr)
    add_text(s2, icon, cc, cy+0.08, cw, 0.4, size=20, align=PP_ALIGN.CENTER)
    add_text(s2, label, cc, cy+0.52, cw, 0.55, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

# Flow diagram at bottom
add_text(s2, "HOW JAGGU WORKS", 0.4, 4.82, 12, 0.3, size=10, bold=True,
         color=GREY_MID, align=PP_ALIGN.CENTER)
arrow_flow(s2, ["Diagnose", "Recommend", "Optimize", "Deliver"],
           l=1.5, t=5.18, w_total=10.3, h=0.62)

velocity_watermark(s2)
add_notes(s2, """SPEAKER NOTES — Slide 2: Introducing JAGGU AI

Position clearly: "JAGGU is not another dashboard. It's an AI consultant that thinks like your best KAM."

Key differentiators:
- JAGGU analyses seller data and gives Observation → Root Cause → Recommendation → Expected Impact
- It covers 6 intelligence layers: Delivery, Courier, Product, Pincode, VAS, Operations
- Built on Velocity GDI — Growth & Delivery Intelligence Engine
- Every answer is backed by the seller's actual shipment data

Demo moment: "If you ask JAGGU why RTO is high, it tells you exactly which courier, which state, which product — and what to do about it TODAY." """)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW JAGGU WORKS
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
gradient_bg(s3)
accent_bar(s3, y=0.0, h=0.06, color=PURPLE_MID)
accent_bar(s3, y=7.44, h=0.06, color=PURPLE_BRIGHT)
add_text(s3, "03 / 05", 12.6, 0.12, 0.6, 0.3, size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)

add_text(s3, "How JAGGU Works", 0.5, 0.15, 8, 0.65, size=36, bold=True, color=WHITE)
add_text(s3, "From raw shipment data to actionable intelligence in seconds",
         0.5, 0.78, 9, 0.35, size=14, color=PURPLE_LIGHT, italic=True)

# Vertical workflow (left side)
card(s3, 0.4, 1.25, 2.8, 5.55, fill=RGBColor(0x0D,0x13,0x22), border=PURPLE_MID)
add_text(s3, "DATA PIPELINE", 0.55, 1.35, 2.5, 0.3, size=10, bold=True, color=GREY_MID)

pipeline = [
    ("📂", "Seller Data",           WHITE,        RGBColor(0x1F,0x29,0x37)),
    ("⚡", "Velocity GDI Engine",   PURPLE_LIGHT, RGBColor(0x1A,0x12,0x2E)),
    ("🤖", "JAGGU AI",              ACCENT_GREEN, RGBColor(0x06,0x20,0x17)),
    ("🔍", "Root Cause Analysis",   ACCENT_AMBER, RGBColor(0x2D,0x1F,0x07)),
    ("✅", "Recommendations",       ACCENT_GREEN, RGBColor(0x06,0x20,0x17)),
    ("📈", "Improved Delivery %",   WHITE,        RGBColor(0x0D,0x13,0x22)),
]
py = 1.75
for icon, label, tc, fc in pipeline:
    add_rect(s3, 0.6, py, 2.4, 0.44, fill_rgb=fc, line_rgb=INDIGO_LINE, line_width=0.4)
    add_text(s3, f"{icon}  {label}", 0.65, py+0.07, 2.3, 0.32, size=11, bold=True, color=tc)
    if label != "📈  Improved Delivery %":
        add_text(s3, "↓", 0.6, py+0.44, 2.4, 0.2, size=9, color=GREY_MID, align=PP_ALIGN.CENTER)
        py += 0.65

# Right side — 5 intelligence modules
mx = 3.45; my_start = 1.25; mw = 4.6; mh = 1.02; gap = 0.15

modules = [
    ("❤️  Seller Health Score",
     "Overall delivery health with benchmark comparison · Health Score 0–100",
     ACCENT_GREEN),
    ("🚚  Courier Intelligence",
     "Delivery rate by courier · Allocation recommendations · Concentration risk",
     PURPLE_BRIGHT),
    ("📦  Product Intelligence",
     "Top selling products · Highest RTO products · Revenue at risk by SKU",
     ACCENT_AMBER),
    ("📍  Pincode Intelligence",
     "Opportunity pincodes · Risk pincodes · High-NDR clusters · Courier mismatch",
     ACCENT_RED),
    ("⚙️  VAS Recommendations",
     "AI Calling · WhatsApp NDR · Order Confirmation · Courier Optimization",
     PURPLE_LIGHT),
]
my = my_start
for title, desc, clr in modules:
    card(s3, mx, my, mw, mh, fill=RGBColor(0x11,0x18,0x27), border=clr)
    add_text(s3, title, mx+0.12, my+0.08, mw-0.2, 0.35, size=12, bold=True, color=WHITE)
    add_text(s3, desc, mx+0.12, my+0.44, mw-0.2, 0.52, size=10, color=GREY_LIGHT)
    my += mh + gap

# Far right — response format
rx = 8.25
card(s3, rx, 1.25, 4.7, 5.55, fill=RGBColor(0x0A,0x0F,0x1E), border=PURPLE_MID)
add_text(s3, "JAGGU RESPONSE FORMAT", rx+0.15, 1.35, 4.4, 0.3, size=10, bold=True, color=GREY_MID)

fmt = [
    ("📊 Observation",    "What the data shows — with real numbers", PURPLE_LIGHT),
    ("🔍 Root Cause",     "Why it's happening — courier, state, product, COD", ACCENT_AMBER),
    ("✅ Recommendation", "Specific Velocity action — not generic advice", ACCENT_GREEN),
    ("📈 Expected Impact","Recoverable shipments or delivery gain", WHITE),
]
fy2 = 1.75
for icon_lbl, desc, clr in fmt:
    add_rect(s3, rx+0.15, fy2, 4.35, 0.88, fill_rgb=RGBColor(0x15,0x1C,0x2E),
             line_rgb=INDIGO_LINE, line_width=0.3)
    add_text(s3, icon_lbl, rx+0.25, fy2+0.05, 4.1, 0.3, size=11, bold=True, color=clr)
    add_text(s3, desc, rx+0.25, fy2+0.38, 4.1, 0.44, size=10, color=GREY_LIGHT)
    fy2 += 0.98 + 0.13

velocity_watermark(s3)
add_notes(s3, """SPEAKER NOTES — Slide 3: How JAGGU Works

Walk through the pipeline left to right:
"Data flows in from Google Sheets or Metabase → Velocity GDI processes it → JAGGU analyses and responds."

Highlight the response format:
"Every JAGGU answer has four parts: Observation (what the data shows), Root Cause (why), Recommendation (what to do), and Expected Impact (how many shipments recovered)."

This is NOT generic AI. Every answer is grounded in the seller's actual data for the selected date range.

Demo cue: Show the JaGau AI page answering 'Why is RTO high?' for a specific seller. Point to Observation → Root Cause → Recommendation → Impact structure. """)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BUSINESS IMPACT
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
gradient_bg(s4)
accent_bar(s4, y=0.0, h=0.06, color=ACCENT_GREEN)
accent_bar(s4, y=7.44, h=0.06, color=PURPLE_MID)
add_text(s4, "04 / 05", 12.6, 0.12, 0.6, 0.3, size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)

add_text(s4, "Business Impact", 0.5, 0.15, 8, 0.65, size=36, bold=True, color=WHITE)
add_text(s4, "What changes when JAGGU is activated",
         0.5, 0.78, 8, 0.35, size=14, color=ACCENT_GREEN, italic=True)

# BEFORE card
card(s4, 0.4, 1.25, 3.7, 5.3, fill=RGBColor(0x1A,0x07,0x07), border=ACCENT_RED)
add_text(s4, "BEFORE  JAGGU", 0.55, 1.35, 3.4, 0.35, size=12, bold=True, color=ACCENT_RED)
add_text(s4, "10,000 Orders / Month", 0.55, 1.75, 3.4, 0.3, size=11, color=GREY_MID)
add_text(s4, "65%", 0.55, 2.08, 1.4, 0.7, size=44, bold=True, color=ACCENT_RED)
add_text(s4, "Delivery\nRate", 1.95, 2.2, 2, 0.5, size=12, color=GREY_LIGHT)
add_rect(s4, 0.6, 2.82, 3.35, 0.02, fill_rgb=INDIGO_LINE)
add_text(s4, "3,500", 0.55, 2.9, 1.4, 0.6, size=36, bold=True, color=ACCENT_RED)
add_text(s4, "Orders\nFailed", 1.95, 2.98, 2, 0.5, size=12, color=GREY_LIGHT)
add_rect(s4, 0.6, 3.55, 3.35, 0.02, fill_rgb=INDIGO_LINE)
add_text(s4, "Reality Today", 0.55, 3.62, 3.4, 0.3, size=10, bold=True, color=GREY_MID)
bullet_block(s4, [
    "No courier intelligence",
    "Manual NDR handling",
    "No pincode restrictions",
    "Reactive operations",
    "Slow decision making",
], l=0.55, t=3.95, w=3.3, size=10, icon="✗", icon_color=ACCENT_RED)

# ARROW
add_text(s4, "→", 4.2, 3.2, 0.6, 0.6, size=28, bold=True, color=PURPLE_LIGHT,
         align=PP_ALIGN.CENTER)
add_text(s4, "JAGGU\nAI", 4.1, 3.78, 0.8, 0.55, size=9, bold=True, color=PURPLE_LIGHT,
         align=PP_ALIGN.CENTER)

# AFTER card
card(s4, 5.0, 1.25, 3.7, 5.3, fill=RGBColor(0x06,0x1C,0x0E), border=ACCENT_GREEN)
add_text(s4, "AFTER  JAGGU", 5.15, 1.35, 3.4, 0.35, size=12, bold=True, color=ACCENT_GREEN)
add_text(s4, "10,000 Orders / Month", 5.15, 1.75, 3.4, 0.3, size=11, color=GREY_MID)
add_text(s4, "82%+", 5.15, 2.08, 1.8, 0.7, size=44, bold=True, color=ACCENT_GREEN)
add_text(s4, "Delivery\nRate", 6.95, 2.2, 1.6, 0.5, size=12, color=GREY_LIGHT)
add_rect(s4, 5.2, 2.82, 3.35, 0.02, fill_rgb=INDIGO_LINE)
add_text(s4, "+1,700", 5.15, 2.9, 1.8, 0.6, size=36, bold=True, color=ACCENT_GREEN)
add_text(s4, "Extra\nDeliveries", 6.95, 2.98, 1.6, 0.5, size=12, color=GREY_LIGHT)
add_rect(s4, 5.2, 3.55, 3.35, 0.02, fill_rgb=INDIGO_LINE)
add_text(s4, "With JAGGU Active", 5.15, 3.62, 3.4, 0.3, size=10, bold=True, color=GREY_MID)
bullet_block(s4, [
    "AI Calling → 38% NDR recovery",
    "WhatsApp NDR → 8% COD RTO save",
    "Courier optimization per state",
    "Pincode COD restrictions",
    "Order Confirmation pre-dispatch",
], l=5.15, t=3.95, w=3.3, size=10, icon="✓", icon_color=ACCENT_GREEN)

# Right — VAS benefits
card(s4, 8.9, 1.25, 4.05, 5.3, fill=RGBColor(0x11,0x18,0x27), border=PURPLE_BRIGHT)
add_text(s4, "JAGGU ACTIVATES", 9.05, 1.35, 3.8, 0.3, size=10, bold=True, color=GREY_MID)

vas_items = [
    ("📞 AI Calling",        "38% NDR recovery rate",    PURPLE_BRIGHT),
    ("💬 WhatsApp NDR",     "8% COD RTO reduction",     ACCENT_GREEN),
    ("🔍 Order Confirm",    "12% pre-dispatch RTO block",ACCENT_AMBER),
    ("🚚 Courier Optimize", "State-level rebalancing",   PURPLE_LIGHT),
    ("📍 Pincode Block",    "Zero-cost COD restriction", ACCENT_RED),
    ("🚀 NDD Activation",  "Zone A/B next-day delivery",ACCENT_GREEN),
]
vy = 1.75
for icon_lbl, desc, clr in vas_items:
    add_rect(s4, 9.05, vy, 3.75, 0.68, fill_rgb=RGBColor(0x0D,0x13,0x22),
             line_rgb=clr, line_width=0.4)
    add_text(s4, icon_lbl, 9.15, vy+0.05, 3.55, 0.28, size=11, bold=True, color=clr)
    add_text(s4, desc, 9.15, vy+0.35, 3.55, 0.28, size=10, color=GREY_LIGHT)
    vy += 0.75 + 0.07

velocity_watermark(s4)
add_notes(s4, """SPEAKER NOTES — Slide 4: Business Impact

Lead with the numbers: "For a 10,000 order/month seller at 65% delivery — JAGGU can add 1,700+ deliveries per month."

Walk through the VAS stack:
- AI Calling: 38% of stuck NDRs recovered via IVR outreach
- WhatsApp NDR: COD buyers re-engaged, 8% fewer RTOs
- Order Confirmation: fake/impulsive orders caught BEFORE dispatch
- Courier Optimization: data shows which courier to use per state
- Pincode Blocking: zero-cost policy change prevents bad-pincode COD

Important: "These are real benchmarks, not projections. JAGGU calculates these from the seller's actual data."

If challenged on numbers: "We show recoverable shipments, not assumed revenue. All calculations are based on actual NDR counts and industry-standard recovery rates." """)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WHY JAGGU WINS
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
gradient_bg(s5)
# Hero gradient header
add_rect(s5, 0, 0.0, 13.33, 2.1, fill_rgb=RGBColor(0x0F,0x0A,0x1E))
accent_bar(s5, y=0.0, h=0.06, color=PURPLE_BRIGHT)
accent_bar(s5, y=2.1, h=0.04, color=PURPLE_MID)
add_text(s5, "05 / 05", 12.6, 0.12, 0.6, 0.3, size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)

add_text(s5, "Why JAGGU Wins", 0.5, 0.18, 8, 0.65, size=38, bold=True, color=WHITE)
add_text(s5, "The shift from reactive dashboards to proactive AI consulting",
         0.5, 0.82, 9, 0.38, size=14, color=PURPLE_LIGHT, italic=True)
add_text(s5, "Hackathon Demo Ready · Production Architecture · Velocity Integrated",
         0.5, 1.42, 9, 0.32, size=11, color=GREY_MID)

# Comparison columns
col_w = 4.15

# Manual / before
card(s5, 0.4, 2.22, col_w, 3.05, fill=RGBColor(0x1A,0x07,0x07), border=ACCENT_RED)
add_text(s5, "❌  Without JAGGU", 0.6, 2.32, 3.9, 0.38, size=13, bold=True, color=ACCENT_RED)
bullet_block(s5, [
    "Manual analysis in Excel/Sheets",
    "Multiple disconnected dashboards",
    "Reactive — act after damage done",
    "Slow decision making (days)",
    "No courier intelligence",
    "No VAS adoption guidance",
    "Generic advice, not seller-specific",
], l=0.6, t=2.75, w=3.8, size=11, icon="✗", icon_color=ACCENT_RED)

# With JAGGU
card(s5, 4.7, 2.22, col_w, 3.05, fill=RGBColor(0x06,0x1C,0x0E), border=ACCENT_GREEN)
add_text(s5, "✅  With JAGGU AI", 4.9, 2.32, 3.9, 0.38, size=13, bold=True, color=ACCENT_GREEN)
bullet_block(s5, [
    "AI KAM — instant consultant",
    "One unified intelligence engine",
    "Proactive — act before damage",
    "Instant decisions (seconds)",
    "Courier optimization per state",
    "VAS recommendations ranked by impact",
    "Personalised per seller, per SKU",
], l=4.9, t=2.75, w=3.8, size=11, icon="✓", icon_color=ACCENT_GREEN)

# Closing statement card
card(s5, 9.05, 2.22, 3.9, 3.05, fill=RGBColor(0x0F,0x0A,0x2E), border=PURPLE_BRIGHT)
add_text(s5, "One AI Consultant", 9.2, 2.38, 3.6, 0.45,
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s5, "Thousands of\nSellers", 9.2, 2.88, 3.6, 0.65,
         size=22, bold=True, color=PURPLE_LIGHT, align=PP_ALIGN.CENTER)
add_text(s5, "Millions of\nDeliveries\nOptimized", 9.2, 3.58, 3.6, 0.82,
         size=18, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

# Final tagline
card(s5, 0.4, 5.45, 12.5, 1.7, fill=RGBColor(0x0F,0x0A,0x2E), border=PURPLE_BRIGHT)
add_text(s5, "JAGGU AI", 0.6, 5.55, 12.0, 0.68,
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s5, "Turning Delivery Data into Growth Decisions",
         0.6, 6.18, 12.0, 0.38, size=16, color=PURPLE_LIGHT,
         italic=True, align=PP_ALIGN.CENTER)
add_text(s5, "Powered by Velocity Shipping  ·  Built with Claude AI",
         0.6, 6.62, 12.0, 0.3, size=10, color=GREY_MID, align=PP_ALIGN.CENTER)

velocity_watermark(s5)
add_notes(s5, """SPEAKER NOTES — Slide 5: Why JAGGU Wins

This is your closing. Own it.

"Dashboards tell you what happened. JAGGU tells you what to do about it — right now, for this seller, based on their actual data."

Three-line pitch for judges:
"One AI consultant. Thousands of sellers. Millions of deliveries optimized."

Why it wins vs competitors:
1. Not a generic chatbot — every answer is grounded in the seller's shipment data
2. Not a dashboard — proactive, conversational, actionable
3. Velocity-native — VAS recommendations are specific to Velocity's product suite
4. Demo-ready — live on Google Sheets data, works in the browser right now

Close with: "JAGGU is live, it's working, and it's ready to scale across Velocity's seller network." """)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
pptx_path = os.path.join(out_dir, "JAGGU_AI_Hackathon_Presentation.pptx")
prs.save(pptx_path)
print(f"✅ PPTX saved: {pptx_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size:   {os.path.getsize(pptx_path):,} bytes")
